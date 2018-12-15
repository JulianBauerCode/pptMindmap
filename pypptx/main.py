import pandas as pd
import collections

#####################################
# Definitions

class Interaction(object):
    # def __init__(self, nameProject):
        # self.nameProject = nameProject
        # self.labelTo = None
        # self.labelFrom = None
        # self.labelWith = None
        
    def __init__(self):
        self.labelTo = None
        self.labelFrom = None
        self.labelWith = None
        
    def addTo(self, label):
        self.labelTo = label

    def addFrom(self, label):
        self.labelFrom = label

    def addWith(self, label):
        self.labelWith = label
    
    def __repr__(self,):
        repr = ''
        for key, value in self.__dict__.items():
            if value:
                repr = repr + '{}:{}  '.format(key, eval('self.'+key))
        return repr

#####################################
# Parameters

blank = 'x'

#####################################
# Read data

supportsExcel = pd.read_excel('supports.xlsx', index_col=0)
cooperationsExcel = pd.read_excel('cooperations.xlsx', index_col=0)

# Check consistency
assert supportsExcel.columns.tolist() == cooperationsExcel.columns.tolist()
assert supportsExcel.index.values.tolist() == supportsExcel.columns.tolist()
assert cooperationsExcel.index.values.tolist() == cooperationsExcel.columns.tolist()

projects = supportsExcel.columns.tolist()

#####################################
# Translate input data into comfortable format

interactions = {}
for project in projects:
    interactions[project] = collections.defaultdict(Interaction)

import collections
dict_x = collections.defaultdict(list)

# supporting interactions
for rowTo in projects:
    for colFrom in projects:
        label = supportsExcel.loc[rowTo][colFrom]
        if label != blank:
            interactions[rowTo][colFrom].addFrom(label)
            interactions[colFrom][rowTo].addTo(label)

# supporting cooperations
for rowTo in projects:
    for colFrom in projects:
        label = cooperationsExcel.loc[rowTo][colFrom]
        if label != blank:
            interactions[rowTo][colFrom].addWith(label)
            interactions[colFrom][rowTo].addWith(label)

#####################################
# Draw graphic

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import MSO_SHAPE

import numpy as np

def xyCentered(x, y, wOffset, hOffset, xOffset=125, yOffset=95):
    return(x + xOffset - wOffset/2.0, -y + yOffset - hOffset/2.0)

def xyByPolar(radius, angle):
    angle = np.deg2rad(angle)
    return (radius * np.cos(angle), radius * np.sin(angle))

radius = 60

width = 20
height = width

project = projects[0]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
xShifted, yShifted = xyCentered(0,0, width, height)
txBox = slide.shapes.add_textbox(Mm(xShifted), Mm(yShifted), Mm(width), Mm(height))
tf = txBox.text_frame
tf.text = project

cInteractions = interactions[project]
nbrPartners = len(cInteractions)
angleIncr = 360.0 / nbrPartners
angle = 0

shapes = slide.shapes

for partner, labels in cInteractions.items():
    print(partner)
    x,y = xyByPolar(radius, angle)
    print(x,y)
    xShifted, yShifted = xyCentered(x,y, width, height)
    print(xShifted, yShifted)
    
    txBox = slide.shapes.add_textbox(Mm(xShifted), Mm(yShifted), Mm(width), Mm(height))
    tf = txBox.text_frame
    tf.text = partner

    for s in [MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.QUAD_ARROW, MSO_SHAPE.LEFT_RIGHT_ARROW,]:
        shape = shapes.add_shape(s, Mm(xShifted), Mm(yShifted), Mm(width), Mm(height))
        shape.text = 'Step 1'


    angle = angle + angleIncr




slide = prs.slides.add_slide(blank_slide_layout)
# left = top = width = height = Cm(1)
txBox = slide.shapes.add_textbox(x,y, width, height)
tf = txBox.text_frame

tf.text = "Here"

# p = tf.add_paragraph()
# p.text = "This is a second paragraph that's bold"
# p.font.bold = True

# p = tf.add_paragraph()
# p.text = "This is a third paragraph that's big"
# p.font.size = Pt(40)

prs.save('test.pptx')